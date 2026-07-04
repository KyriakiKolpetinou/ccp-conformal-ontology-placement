#!/usr/bin/env python
"""Build a concise findings deck for the Calibrated Concept Placement project (run with base python)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "phase2", "figures")
BLUE = RGBColor(0x1F, 0x4E, 0x79); GREY = RGBColor(0x59, 0x59, 0x59)
RED = RGBColor(0xB0, 0x2418 // 256 if False else 0x24, 0x18); GREEN = RGBColor(0x2C, 0x7A, 0x2C)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True; return box.text_frame

def para(tf, text, size=18, bold=False, color=None, bullet=False, space=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    run = p.add_run(); run.text = ("•  " if bullet else "") + text
    run.font.size = Pt(size); run.font.bold = bold
    if color: run.font.color.rgb = color
    return p

def title_slide(title, sub):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.5), W, Inches(1.5)); bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = tb(s, 0.6, 2.6, 12.1, 1.3); p = para(tf, title, 34, True, RGBColor(255,255,255), first=True); p.alignment = PP_ALIGN.CENTER
    tf2 = tb(s, 0.6, 4.3, 12.1, 1.2); p = para(tf2, sub, 18, False, GREY, first=True); p.alignment = PP_ALIGN.CENTER
    return s

def header(s, text):
    tf = tb(s, 0.5, 0.3, 12.3, 0.9); para(tf, text, 26, True, BLUE, first=True)

def content_slide(title, bullets):
    s = prs.slides.add_slide(BLANK); header(s, title)
    tf = tb(s, 0.7, 1.4, 12.0, 5.6)
    for i, b in enumerate(bullets):
        if isinstance(b, tuple): txt, sz, col, bul = b
        else: txt, sz, col, bul = b, 18, None, True
        para(tf, txt, sz, False, col, bul, first=(i == 0))
    return s

def figure_slide(title, fig, caption, bullets):
    s = prs.slides.add_slide(BLANK); header(s, title)
    s.shapes.add_picture(os.path.join(FIG, fig), Inches(0.5), Inches(1.5), height=Inches(5.3))
    tf = tb(s, 7.4, 1.5, 5.5, 5.3)
    para(tf, caption, 16, True, BLUE, first=True)
    for b in bullets: para(tf, b, 16, False, None, True)
    return s

# 1 title
title_slide("Calibrated Concept Placement", "Hierarchy-aware conformal prediction for ontology concept placement\nFindings as of 2026-06-25  ·  MM-S14-Disease (CPP benchmark running)")

# 2 problem + bottom line
content_slide("The task, and the bottom line", [
    ("TASK: a clinical text mentions a NEW disease not yet in SNOMED CT. Place it by choosing the right", 18, None, False),
    ("(parent → child) insertion edge among ~238,000 candidates. Analogy: file it in the right drawer of a huge cabinet.", 18, None, False),
    ("", 8, None, False),
    ("BOTTOM LINE:", 19, BLUE, False),
    ("Guaranteeing the EXACT edge is impossible here — but guaranteeing the correct REGION is feasible and", 18, GREEN, True),
    ("useful to a curator. We wrap the placement model in a hierarchy-aware conformal layer to do this.", 18, GREEN, False),
    ("", 8, None, False),
    ("Backbone: Oxford LM concept-placement model (ESWC 2024). Our contribution = the calibration layer.", 16, GREY, True),
])

# 3 Phase 0
content_slide("Phase 0 — reproduction holds (GO)", [
    "Rebuilt their full pipeline (Edge-Bi-encoder + Edge-Cross-encoder, top-50). Faithful to the paper.",
    ("Our test InR_any@10 = 20.3% vs paper 26.5%  (within their reported validation/test variance).", 18, None, True),
    ("Retrieval recall@50 = 49.6% vs paper 50.0%  — near-exact match on the retrieval step.", 18, None, True),
    ("Several metrics match or exceed theirs (e.g. InR_all@10 10.1% vs 8.7%).", 18, None, True),
    ("→ Solid, comparable foundation to build the calibration layer on.", 18, GREEN, True),
])

# 4 Phase 1 walls
content_slide("Phase 1 — naive calibration hits two walls", [
    ("WALL 1: exact-edge prediction sets are not viable.", 19, BLUE, False),
    "Coverage is hard-capped at recall@pool (~50%): the gold edge is simply not retrieved half the time.",
    "To promise >50% coverage, the set explodes to the entire 50-edge pool — useless.",
    ("WALL 2: the model's confidence is anti-correlated with exact correctness.", 19, BLUE, False),
    "Its MOST-confident predictions (score ≈ 0.999) are correct 0% of the time — 'confident false positives'.",
    ("→ Both naive fallbacks (small sets, abstain-by-confidence) fail. This is the wall the project anticipated.", 18, RED, True),
])

# 5 the insight
content_slide("The insight that rescues it: near-misses are NEAR", [
    ("When the model misses the EXACT edge, its prediction is still ontologically CLOSE to the gold.", 19, GREEN, False),
    "Measured with Wu–Palmer similarity over the SNOMED hierarchy (parent + child placement):",
    ("Exact-edge coverage (top-10): 0.20", 18, None, True),
    ("Right-region coverage (full-edge WP≥0.7): 0.77   ·   close (WP≥0.8): 0.62", 18, None, True),
    ("Among exact MISSES: average closeness ≈ 0.81 (about one hop from gold).", 18, None, True),
    ("Filing-cabinet: it rarely picks the exact folder, but almost always opens the right drawer.", 18, GREY, True),
    ("This is the 'lenient evaluation' the original authors left as future work — we calibrate on it.", 16, GREY, True),
])

# 6 fig1 risk-coverage
figure_slide("Result 1 — confidence tracks REGION, not exact edge", "fig1_risk_coverage.png",
    "Selective prediction (auto-accept the most confident):", [
    "Region quality (blue) stays 0.78–0.91.",
    "Exact correctness (red) hugs zero — and is 0 at the highest confidence.",
    "Continuous risk AURC: 0.15 (region) vs 0.96 (exact).",
    "→ The model knows when it's in the right region; usable for curator triage.",
])

# 7 fig3 region sets
figure_slide("Result 2 — small coverage-guaranteed region sets", "fig3_setsize_coverage.png",
    "How many edges a curator reviews vs coverage:", [
    "WP≥0.7: ~10-edge set → 65% coverage; ~30 → 88%.",
    "WP≥0.8 (closer): ~19 edges → 65%.",
    "Exact / WP≥0.9: only exist at 27–50 edges (no small set).",
    "→ Relaxing exact→region makes small, trustworthy sets possible.",
])

# 8 calibration + mondrian
figure_slide("Calibration & subgroup fairness (with an honest limit)", "fig5_mondrian_depth.png",
    "Class-conditional (Mondrian) conformal:", [
    "Marginal conformal HIDES subgroup gaps (mid-depth 57% vs shallow 100%).",
    "Mondrian rebalances: worst group improves, easy groups get tighter sets.",
    "BUT it can't fix the overall undershoot — test concepts are uniformly harder than calibration ones (concept drift).",
    "Honest limitation the paper must own (see also fig2 calibration, fig4 drift).",
])

# 9 status + next
content_slide("Status, publishability, next steps", [
    ("STATUS:  reproduction GO  ·  exact-edge calibration NO-GO  ·  region calibration WORKS.", 18, GREEN, True),
    ("PUBLISHABILITY (honest): a plausible mid-tier paper (JBI / JAMIA Open / J. Biomed. Semantics).", 18, None, True),
    "   Strength: novel reliability/curator-usability angle; the confidence-tracks-region finding.",
    "   Risks to manage: the coverage guarantee undershoots (drift); single small dataset; modest novelty.",
    "   → Reframe headline to risk-coverage/curator-triage (robust to drift), and add a 2nd dataset.",
    ("NEXT:  CPP benchmark running now (job 326, 432 test mentions) → cross-dataset generalization;", 18, BLUE, True),
    "          then class-conditional refinement + write-up.",
])

out = os.path.join(HERE, "CCP_findings.pptx"); prs.save(out)
print("saved:", out, "(%d slides)" % len(prs.slides._sldIdLst))
